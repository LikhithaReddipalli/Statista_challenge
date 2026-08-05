"""Deliberately malformed inputs, exercised end-to-end.

Runs the real pipeline functions against real malformed fixtures to prove
graceful degradation actually happens. Not pytest-collected (no bare
asserts, no Test* classes) — run directly: python tests/test_edge_cases.py
"""

import io
import json
import os
import sys
from contextlib import contextmanager
from pathlib import Path
from typing import Any, Optional

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from pptx import Presentation
from pptx.util import Emu
from pydantic import ValidationError

from main import generate_presentation
from models import Metric, SlideSpec
from planner import _chart_type_for_metrics, _parse_and_validate
from renderer import build_presentation
from validator import build_context, normalize_metrics, validate_input

FIXTURES = Path(__file__).parent / "malformed_inputs"
SAMPLE_DATA = Path(__file__).resolve().parent.parent / "sample_data" / "gen_z_purchase_behavior_analysis.json"


class _Counter:
    """Pass/fail tally for the run (instance, not a module global, so
    importing individual tests elsewhere can't leak shared state)."""

    def __init__(self) -> None:
        self.failures = 0

    def check(self, label: str, condition: bool) -> None:
        status = "PASS" if condition else "FAIL"
        print(f"[{status}] {label}")
        if not condition:
            self.failures += 1


counter = _Counter()
check = counter.check


@contextmanager
def _env(**kv: Optional[str]):
    """Temporarily set env vars, restoring prior values on exit. A value
    of None unsets that key for the duration of the block."""
    original = {k: os.environ.get(k) for k in kv}
    for k, v in kv.items():
        if v is None:
            os.environ.pop(k, None)
        else:
            os.environ[k] = v
    try:
        yield
    finally:
        for k, v in original.items():
            if v is None:
                os.environ.pop(k, None)
            else:
                os.environ[k] = v


@contextmanager
def _patched(obj: Any, attr: str, value: Any):
    """Temporarily replace obj.attr, restoring the original on exit."""
    original = getattr(obj, attr)
    setattr(obj, attr, value)
    try:
        yield
    finally:
        setattr(obj, attr, original)


def _max_text_bottom(slide) -> int:
    """Furthest-down edge (in EMU) of any non-empty text shape on a slide."""
    return max(
        (shape.top + shape.height for shape in slide.shapes if shape.has_text_frame and shape.text_frame.text.strip()),
        default=Emu(0),
    )


def _slide_spec_stub(**overrides: Any) -> SlideSpec:
    """Build a SlideSpec via model_construct (bypassing __init__
    validation) with sane defaults, so a test only spells out the field(s)
    it's deliberately breaking. Simulates a construction path that skips
    validation — the renderer must still not crash on it."""
    defaults: dict[str, Any] = dict(
        title="x",
        subtitle=None,
        layout="chart",
        chart_type=None,
        key_message=None,
        bullets=[],
        metrics=[],
        kpis=[],
        source=None,
        notes=None,
    )
    defaults.update(overrides)
    return SlideSpec.model_construct(**defaults)


def test_empty_key_insights() -> None:
    print("\n--- empty key_insights: should NOT raise, should produce a summary-only deck ---")
    data = json.loads((FIXTURES / "empty_key_insights.json").read_text())

    analysis = validate_input(data)
    check("validate_input does not raise", True)
    check("validate_input flags _summary_only", analysis.get("_summary_only") is True)

    result = generate_presentation(data, use_llm=False)
    check("generate_presentation succeeds", result.slide_count > 0)
    check("summary-only deck has exactly 2 slides (title + summary)", result.slide_count == 2)


def test_missing_source() -> None:
    print("\n--- missing source: renderer should show 'Source unavailable.' not fail ---")
    data = json.loads((FIXTURES / "missing_source.json").read_text())

    result = generate_presentation(data, use_llm=False)
    check("generate_presentation succeeds despite missing source", result.slide_count > 0)

    prs = Presentation(io.BytesIO(result.pptx_bytes))
    chart_slide = prs.slides[1]  # slide 0 = title, slide 1 = the one insight
    footer_texts = [
        shape.text_frame.text
        for shape in chart_slide.shapes
        if shape.has_text_frame and "Source" in shape.text_frame.text
    ]
    check(
        "chart slide footer reads 'Source unavailable.'",
        any("Source unavailable." in t for t in footer_texts),
    )


def test_llm_response_missing_field() -> None:
    print("\n--- LLM response missing a required field: planner must return None (fallback trigger) ---")
    raw_text = (FIXTURES / "llm_response_missing_field.json").read_text()

    result = _parse_and_validate(raw_text)
    check("_parse_and_validate returns None on missing required 'title'", result is None)


def test_llm_response_broken_json_syntax() -> None:
    print("\n--- LLM response is syntactically invalid JSON: planner must return None (fallback trigger), not raise ---")
    raw_text = (FIXTURES / "broken_syntax.json").read_text()

    result = _parse_and_validate(raw_text)
    check("_parse_and_validate returns None on invalid JSON syntax", result is None)


def test_missing_analysis_key() -> None:
    print("\n--- JSON with neither 'analysis', 'title', nor 'key_insights': must raise a descriptive ValueError ---")
    try:
        validate_input({"not_analysis": {}})
        check("validate_input raises ValueError", False)
    except ValueError as e:
        check("validate_input raises ValueError", True)
        check("error message mentions 'analysis'", "analysis" in str(e))


def test_unwrapped_analysis_object() -> None:
    print("\n--- JSON with no 'analysis' wrapper key: top-level object should be used as-is ---")
    data = {"title": "No Wrapper", "key_insights": []}

    analysis = validate_input(data)
    check("validate_input does not raise on an unwrapped object", True)
    check("validate_input returns the top-level dict itself", analysis is data)
    check("validate_input still flags _summary_only", analysis.get("_summary_only") is True)

    result = generate_presentation(data, use_llm=False)
    check("generate_presentation succeeds on an unwrapped object", result.slide_count == 2)


def test_summary_slide_bullet_overflow() -> None:
    print("\n--- many insight categories: summary slide must not run content off-slide ---")
    data = json.loads(SAMPLE_DATA.read_text())
    base_insights = data["analysis"]["key_insights"]
    data["analysis"]["key_insights"] = base_insights * 3  # 9 bullets — more than fit

    result = generate_presentation(data, use_llm=False)

    prs = Presentation(io.BytesIO(result.pptx_bytes))
    slide_height_emu = prs.slide_height
    summary_slide = prs.slides[-1]

    check("no summary-slide text renders past the bottom of the slide", _max_text_bottom(summary_slide) <= slide_height_emu)

    all_text = " ".join(
        shape.text_frame.text for shape in summary_slide.shapes if shape.has_text_frame
    )
    check("excess bullets are silently dropped, not shown as a '+N more' label", "more" not in all_text.lower())


def test_text_only_slide_bullet_overflow() -> None:
    print("\n--- text-only chart slide with many bullets: must not run off-slide either ---")
    data = {
        "analysis": {
            "title": "Overflow Check",
            "key_insights": [
                {
                    "category": "Everything At Once",
                    "insight": "A slide with far more bullets than a slide can hold.",
                    "metrics": [],  # no metrics -> text-only body
                    "source": {"name": "Test Fixture"},
                }
            ],
        }
    }
    # Patch in a long bullet list the way an LLM-planned slide might.
    from planner import plan_rule_based

    analysis = validate_input(data)
    insights = analysis.get("key_insights", [])
    normalized = normalize_metrics(insights)
    context = build_context(analysis, normalized)
    slides = plan_rule_based(context)
    slides[1].bullets = [f"Bullet point number {i}" for i in range(20)]

    buf = io.BytesIO()
    build_presentation(slides, buf)
    prs = Presentation(io.BytesIO(buf.getvalue()))
    slide_height_emu = prs.slide_height
    chart_slide = prs.slides[1]

    check("no text-only-body content renders past the bottom of the slide", _max_text_bottom(chart_slide) <= slide_height_emu)


def test_pie_chart_type_for_part_to_whole_metrics() -> None:
    print("\n--- rule-based planner: percentages summing to ~100 become a pie, not a bar ---")
    part_to_whole = [
        Metric(label="North America", value=45, kind="comparison"),
        Metric(label="Europe", value=35, kind="comparison"),
        Metric(label="Asia-Pacific", value=20, kind="comparison"),
    ]
    check("3 metrics summing to 100 -> pie", _chart_type_for_metrics(part_to_whole) == "pie")

    independent_stats = [
        Metric(label="Influencer posts", value=57, kind="comparison"),
        Metric(label="Brand social posts", value=54, kind="comparison"),
        Metric(label="Feed advertisements", value=49, kind="comparison"),
        Metric(label="Friend recommendations", value=40, kind="comparison"),
    ]
    check(
        "4 independent stats not summing to 100 -> stays comparison bar",
        _chart_type_for_metrics(independent_stats) == "comparison",
    )

    too_many_slices = [Metric(label=f"Segment {i}", value=100 / 8, kind="comparison") for i in range(8)]
    check(
        "8 slices summing to 100 -> still bar (too many slices for a readable pie)",
        _chart_type_for_metrics(too_many_slices) == "comparison",
    )


def test_chart_type_without_metrics() -> None:
    print("\n--- chart_type set but metrics empty: must not reach the renderer's add_chart crash ---")

    for chart_type in ("ranking", "comparison", "trend", "pie"):
        try:
            SlideSpec(title="x", layout="chart", chart_type=chart_type, metrics=[])
            check(f"SlideSpec rejects chart_type={chart_type!r} with empty metrics", False)
        except ValidationError:
            check(f"SlideSpec rejects chart_type={chart_type!r} with empty metrics", True)

    # Second line of defense: a SlideSpec built via model_construct
    # (skips validation) must still degrade to text-only, not crash.
    spec = _slide_spec_stub(chart_type="ranking", metrics=[])
    buf = io.BytesIO()
    try:
        build_presentation([spec], buf)
        check("renderer degrades to text-only instead of crashing on empty-metrics chart_type", True)
    except Exception:
        check("renderer degrades to text-only instead of crashing on empty-metrics chart_type", False)

    prs = Presentation(io.BytesIO(buf.getvalue()))
    all_text = " ".join(
        shape.text_frame.text for shape in prs.slides[0].shapes if shape.has_text_frame
    )
    check("degraded slide still shows the title as placeholder content", "x" in all_text)


def test_llm_planned_deck_render_failure_falls_back() -> None:
    print("\n--- plan_llm returns a slide that crashes build_presentation: generate_presentation must fall back, not crash ---")
    import main

    # A shape model_validator can't catch (built via model_construct) that
    # still breaks the renderer — proves generate_presentation's try/except
    # is a real second line of defense (see test_chart_type_without_metrics
    # for the model_validator guard itself).
    bad_slides = [_slide_spec_stub(title="Bad Slide", layout="kpi", kpis="not-a-list")]

    with _patched(main, "plan_llm", lambda context, user_instruction=None: bad_slides):
        data = json.loads(SAMPLE_DATA.read_text())
        result = generate_presentation(data, use_llm=True)
        check("generate_presentation does not crash on a bad LLM-planned deck", result.slide_count > 0)
        check("fell_back is set when the LLM-planned deck fails to render", result.fell_back is True)


class _FakeResponse:
    """Stand-in for langchain's AIMessage — plan_llm only reads `.content`."""

    def __init__(self, content):
        self.content = content


class _FakeChatGoogleGenerativeAI:
    """Stand-in for ChatGoogleGenerativeAI that returns a canned response
    from .invoke(), so plan_llm() can run end-to-end without a real network
    call or API key. `response` is always passed explicitly — no shared
    class-level state — so nothing leaks between the sub-tests below."""

    def __init__(self, model=None, google_api_key=None, response=None, **kwargs):
        self.model = model
        self.google_api_key = google_api_key
        self.response = response

    def invoke(self, prompt):
        return self.response


def _make_fake_client_factory(response):
    """Return a client class bound to a fixed canned response, and a mutable
    box to read back the kwargs the real code constructed it with."""
    captured: dict[str, Any] = {}

    class _Client(_FakeChatGoogleGenerativeAI):
        def __init__(self, model=None, google_api_key=None, **kwargs):
            captured["model"] = model
            super().__init__(model=model, google_api_key=google_api_key, response=response, **kwargs)

    return _Client, captured


def test_plan_llm_end_to_end_happy_path() -> None:
    print("\n--- plan_llm end-to-end: mocked Gemini client, real prompt/parse/validate path ---")
    import langchain_google_genai
    import planner

    data = json.loads((FIXTURES / "missing_source.json").read_text())
    analysis = validate_input(data)
    insights = analysis.get("key_insights", [])
    normalized = normalize_metrics(insights)
    context = build_context(analysis, normalized)

    valid_response = json.dumps(
        [
            {
                "title": "Test Deck",
                "layout": "title",
                "subtitle": "Key Insights",
            },
            {
                "title": "Missing Source Example",
                "layout": "chart",
                "chart_type": "comparison",
                "key_message": "Some action leads.",
                "metrics": [
                    {"label": "Some action", "value": 42, "kind": "comparison", "highlight": True},
                    {"label": "Another action", "value": 30, "kind": "comparison"},
                ],
                "source": "Test Fixture",
            },
        ]
    )
    fake_client, captured = _make_fake_client_factory(_FakeResponse(valid_response))

    with _env(GEMINI_API_KEY="test-key-not-real", GEMINI_MODEL="gemini-test-model"), \
         _patched(langchain_google_genai, "ChatGoogleGenerativeAI", fake_client):
        slides = planner.plan_llm(context, user_instruction="Keep it short.")

    check("plan_llm returns a validated list of SlideSpec, not None", slides is not None)
    check("plan_llm returns exactly the 2 slides in the mocked response", slides is not None and len(slides) == 2)
    check(
        "each returned item is a real SlideSpec instance",
        slides is not None and all(isinstance(s, SlideSpec) for s in slides),
    )
    check(
        "GEMINI_MODEL env var reaches the client constructor",
        captured.get("model") == "gemini-test-model",
    )


def test_plan_llm_end_to_end_typed_content_blocks() -> None:
    print("\n--- plan_llm end-to-end: mocked client returns content as typed blocks (thinking-model shape) ---")
    import langchain_google_genai
    import planner

    data = json.loads((FIXTURES / "missing_source.json").read_text())
    analysis = validate_input(data)
    insights = analysis.get("key_insights", [])
    normalized = normalize_metrics(insights)
    context = build_context(analysis, normalized)

    block_response = [
        {"type": "thought_signature", "signature": "ignored-not-text"},
        {"type": "text", "text": '[{"title": "Only Slide", "layout": "title"}]'},
    ]
    fake_client, _captured = _make_fake_client_factory(_FakeResponse(block_response))

    with _env(GEMINI_API_KEY="test-key-not-real"), \
         _patched(langchain_google_genai, "ChatGoogleGenerativeAI", fake_client):
        slides = planner.plan_llm(context)

    check(
        "plan_llm extracts text from typed content blocks and ignores non-text blocks",
        slides is not None and len(slides) == 1 and slides[0].title == "Only Slide",
    )


def test_plan_llm_end_to_end_client_raises() -> None:
    print("\n--- plan_llm end-to-end: mocked client raises (e.g. network/auth error) -> None, no exception ---")
    import langchain_google_genai
    import planner

    data = json.loads((FIXTURES / "missing_source.json").read_text())
    analysis = validate_input(data)
    insights = analysis.get("key_insights", [])
    normalized = normalize_metrics(insights)
    context = build_context(analysis, normalized)

    class _RaisingClient(_FakeChatGoogleGenerativeAI):
        def invoke(self, prompt):
            raise RuntimeError("simulated network failure")

    with _env(GEMINI_API_KEY="test-key-not-real"), \
         _patched(langchain_google_genai, "ChatGoogleGenerativeAI", _RaisingClient):
        slides = planner.plan_llm(context)

    check("plan_llm returns None (not an exception) when the client raises", slides is None)


def test_invalid_json_via_cli() -> None:
    print("\n--- CLI given a non-existent file / invalid JSON: must exit cleanly, not stack-trace ---")
    import subprocess

    # sys.executable (this interpreter) rather than a hardcoded venv path,
    # so this runs under whatever environment actually invoked the suite —
    # a fixed "virtual/bin/python" path breaks on Windows and on any
    # machine/CI image where the venv isn't literally named "virtual".
    # sys.executable, not a hardcoded venv path, so this runs under
    # whatever environment invoked the suite.
    main_py = str(Path(__file__).resolve().parent.parent / "main.py")

    proc = subprocess.run(
        [sys.executable, main_py, "/nonexistent/path.json"],
        capture_output=True,
        text=True,
    )
    check("exits with non-zero status on missing file", proc.returncode != 0)
    check("prints a clean error, not a traceback", "Traceback" not in proc.stderr)


if __name__ == "__main__":
    test_empty_key_insights()
    test_missing_source()
    test_llm_response_missing_field()
    test_llm_response_broken_json_syntax()
    test_missing_analysis_key()
    test_unwrapped_analysis_object()
    test_summary_slide_bullet_overflow()
    test_text_only_slide_bullet_overflow()
    test_pie_chart_type_for_part_to_whole_metrics()
    test_chart_type_without_metrics()
    test_llm_planned_deck_render_failure_falls_back()
    test_plan_llm_end_to_end_happy_path()
    test_plan_llm_end_to_end_typed_content_blocks()
    test_plan_llm_end_to_end_client_raises()
    test_invalid_json_via_cli()

    print(f"\n{'='*50}")
    if counter.failures:
        print(f"{counter.failures} check(s) FAILED")
        sys.exit(1)
    print("All checks passed")
