"""Deterministic renderer: SlideSpec objects -> a .pptx file.

Pure python-pptx, no matplotlib. Charts use python-pptx's native chart API
(CategoryChartData + add_chart) so the result is a real, editable
PowerPoint chart object rather than a static image.

This module only ever receives validated SlideSpec objects (enforced by
main.py), so it never has to defensively handle malformed input.
"""

import logging
from typing import Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Emu, Inches, Pt

from models import SlideSpec

logger = logging.getLogger(__name__)

# --- theme -------------------------------------------------------------
# Statista report palette: one brand hue (deep blue), flat design, no
# shadows or decorative fills. Emphasis is always the same blue, one shade
# darker — never a second color.
PRIMARY = RGBColor(0x00, 0x3A, 0x70)          # deep blue — titles, emphasis
BLUE_MEDIUM = RGBColor(0x6C, 0x8E, 0xBF)      # ordinary chart marks
TEXT_DARK = RGBColor(0x1F, 0x29, 0x37)        # body text
TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)       # sources / footnotes
BACKGROUND = RGBColor(0xFF, 0xFF, 0xFF)       # clean white
GRIDLINE = RGBColor(0xE5, 0xE7, 0xEB)         # light divider / gridlines

# Tints of the brand blue for pie slices (darkest to lightest). Six steps
# covers the pie's 2-6 slice range (see planner._looks_like_part_to_whole);
# a 7th slice would reuse the darkest tint, but pies are capped at 6.
PIE_PALETTE = [
    RGBColor(0x00, 0x3A, 0x70),  # PRIMARY
    RGBColor(0x2C, 0x5E, 0x94),
    RGBColor(0x58, 0x81, 0xB3),
    RGBColor(0x84, 0xA5, 0xCE),
    RGBColor(0xB0, 0xC8, 0xE3),
    RGBColor(0xD6, 0xE3, 0xF0),
]

FONT = "Calibri"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Top of the footer band (source/notes left, page number right) — nothing
# above this line so content never runs under the footer.
FOOTER_TOP = SLIDE_HEIGHT - Inches(1.15)


def _chart_height(top):
    """How tall a chart can be, given where it starts, to stop clear of the footer."""
    return FOOTER_TOP - Inches(0.15) - top


def _add_kpi_card(slide: Slide, left, top, width, height, top_border_height) -> None:
    """A flat white KPI card: square corners, hairline border, no shadow,
    plus a hairline blue top border — the full chrome shared by every
    card in a KPI row (see render_kpi_slide)."""
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = BACKGROUND
    card.line.color.rgb = GRIDLINE
    card.line.width = Pt(0.5)
    card.shadow.inherit = False

    top_border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, top_border_height)
    top_border.fill.solid()
    top_border.fill.fore_color.rgb = BLUE_MEDIUM
    top_border.line.fill.background()
    top_border.shadow.inherit = False


def build_presentation(slides: list[SlideSpec], output) -> None:
    """Render a full deck from validated slides and save it.

    `output` is anything python-pptx's Presentation.save() accepts — a
    file path or a file-like object — so the same function serves both the
    CLI (disk) and the Streamlit UI (in-memory buffer).
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]  # fully blank; we build everything ourselves

    deck_title = slides[0].title if slides else ""
    total = len(slides)

    for i, spec in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        if spec.layout == "title":
            render_title_slide(slide, spec)
        elif spec.layout == "summary":
            render_summary_slide(slide, spec)
        elif spec.layout == "kpi":
            render_kpi_slide(slide, spec)
        else:
            render_chart_slide(slide, spec)

        _add_page_footer(slide, deck_title, i + 1, total)

    prs.save(output)


# --- shared helpers ------------------------------------------------------


def _add_textbox(slide: Slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    box.text_frame.word_wrap = True
    return box


def _style_title(paragraph, size=29, color=PRIMARY, bold=True):
    """Shared title styling — used for both the title slide's headline and
    every other layout's slide header."""
    paragraph.font.name = FONT
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = color


def _add_note(slide: Slide, text: str, left, top, width, height, size=Pt(9)) -> None:
    """Small italic muted "Note: ..." line — shared by the title slide and
    the source footer, the two places a slide's `notes` field is shown."""
    box = _add_textbox(slide, left, top, width, height)
    p = box.text_frame.paragraphs[0]
    p.text = f"Note: {text}"
    p.font.name = FONT
    p.font.size = size
    p.font.italic = True
    p.font.color.rgb = TEXT_MUTED


def _cap_items(items: list[str], max_count: int) -> list[str]:
    """Show at most `max_count` items — anything beyond that is silently
    dropped rather than shown as a "+N more" label, since the label isn't
    clickable/expandable and just adds noise without adding information."""
    return items[:max_count]


def _truncate(text: str, max_len: int) -> str:
    """Hard length cap with an ellipsis for single-line fields with no
    overflow handling of their own (KPI fields, chart labels)."""
    if len(text) <= max_len:
        return text
    return text[: max_len - 1].rstrip() + "…"


def _add_key_message(slide: Slide, text: str, top) -> None:
    """Small italic muted lead-in line above a chart or KPI row."""
    msg_box = _add_textbox(slide, Inches(0.5), top, Inches(12.3), Inches(0.6))
    mp = msg_box.text_frame.paragraphs[0]
    mp.text = text
    mp.font.name = FONT
    mp.font.size = Pt(15)
    mp.font.italic = True
    mp.font.color.rgb = TEXT_MUTED


# Shared baseline for both footer lines so they sit on the same visual row.
PAGE_FOOTER_TOP = SLIDE_HEIGHT - Inches(0.55)


def _add_source_footer(slide: Slide, source: Optional[str], notes: Optional[str]) -> None:
    """Bottom-left source footer, present on every slide. Falls back to an
    explicit "Source unavailable." rather than omitting the footer, so a
    missing source doesn't look like a rendering bug."""
    source_height = Inches(0.5)
    footer_top = PAGE_FOOTER_TOP
    text = source if source else "Source unavailable."
    box = _add_textbox(slide, Inches(0.5), footer_top, Inches(8.3), source_height)
    p = box.text_frame.paragraphs[0]
    p.text = f"Source: {text}"
    p.font.name = FONT
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.color.rgb = TEXT_MUTED

    if notes:
        note_height = Inches(0.35)
        _add_note(slide, notes, Inches(0.5), FOOTER_TOP - note_height, Inches(10.5), note_height)


def _add_page_footer(slide: Slide, deck_title: str, page_number: int, total_pages: int) -> None:
    """Bottom-right footer: deck title + page number, on every slide."""
    box = _add_textbox(slide, Inches(9.3), PAGE_FOOTER_TOP, Inches(3.5), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = f"{deck_title}  |  {page_number} / {total_pages}"
    p.font.name = FONT
    p.font.size = Pt(9)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.RIGHT


def _add_accent_rule(slide: Slide, top) -> None:
    """Thin horizontal accent line under the slide title."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), top, Inches(2.2), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY
    line.line.fill.background()


def _add_slide_header(slide: Slide, spec: SlideSpec) -> None:
    """Title + accent rule, shared by every non-title layout."""
    title_box = _add_textbox(slide, Inches(0.5), Inches(0.35), Inches(11.5), Inches(0.7))
    _style_title(title_box.text_frame.paragraphs[0], size=25)
    title_box.text_frame.paragraphs[0].text = spec.title
    _add_accent_rule(slide, Inches(0.95))


# Where a slide's body starts below the header (chart and KPI slides).
BODY_TOP = Inches(1.2)


# --- title slide -----------------------------------------------------------


def render_title_slide(slide: Slide, spec: SlideSpec) -> None:
    """Cover slide: white background, bold deep-blue title, gray subtitle."""
    title_box = _add_textbox(slide, Inches(1), Inches(2.8), Inches(11.3), Inches(1.8))
    _style_title(title_box.text_frame.paragraphs[0], size=36)
    title_box.text_frame.paragraphs[0].text = spec.title

    if spec.subtitle:
        sub_box = _add_textbox(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(0.7))
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = spec.subtitle
        sp.font.name = FONT
        sp.font.size = Pt(18)
        sp.font.color.rgb = TEXT_DARK

    if spec.notes:
        _add_note(slide, spec.notes, Inches(1), Inches(6.6), Inches(11.3), Inches(0.6), size=Pt(10))


# --- chart slide -------------------------------------------------------

MAX_CHART_BULLETS = 3
CHART_BULLET_LINE_HEIGHT = Pt(14 + 6)
CHART_BULLETS_GAP = Inches(0.1)  # space between chart bottom and the bullets band


def _chart_bullets_height(bullets: list[str]):
    """Vertical space to reserve under a chart for spec.bullets, including
    the gap above the band — zero if there are none."""
    if not bullets:
        return Emu(0)
    rows = min(len(bullets), MAX_CHART_BULLETS)
    return CHART_BULLETS_GAP + CHART_BULLET_LINE_HEIGHT * rows


def _add_chart_bullets(slide: Slide, bullets: list[str], top) -> None:
    """Compact bullet band under a chart, capped to MAX_CHART_BULLETS."""
    display_bullets = _cap_items(bullets, MAX_CHART_BULLETS)

    box = _add_textbox(slide, Inches(0.9), top, Inches(11.5), CHART_BULLET_LINE_HEIGHT * len(display_bullets))
    tf = box.text_frame
    first = True

    for bullet in display_bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = f"•  {bullet}"
        p.font.name = FONT
        p.font.size = Pt(14)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(6)
        first = False


def render_chart_slide(slide: Slide, spec: SlideSpec) -> None:
    """Insight slide: ranking/comparison bar, trend line, pie, or (if no
    chart_type/metrics) a text-only body — never skipped."""
    _add_slide_header(slide, spec)

    body_top = BODY_TOP

    # A chart_type with no metrics can't be plotted; degrade to text-only.
    # SlideSpec's model_validator already rejects this for LLM output —
    # this is a second guard for any other construction path.
    has_chart = spec.chart_type in ("ranking", "comparison", "trend", "pie") and bool(spec.metrics)

    if has_chart:
        if spec.key_message:
            _add_key_message(slide, spec.key_message, body_top)
            chart_top = body_top + Inches(0.7)
        else:
            chart_top = body_top

        # Bullets are opt-in on a chart slide (see planner's prompt rules).
        bullets_height = _chart_bullets_height(spec.bullets)
        chart_height = _chart_height(chart_top) - bullets_height

        if spec.chart_type == "trend":
            _add_line_chart(slide, spec, chart_top, chart_height)
        elif spec.chart_type == "pie":
            _add_pie_chart(slide, spec, chart_top, chart_height)
        else:
            _add_bar_chart(slide, spec, chart_top, chart_height)

        if bullets_height:
            # chart_height already accounts for CHART_BULLETS_GAP, so no
            # extra gap is added here.
            _add_chart_bullets(slide, spec.bullets, chart_top + chart_height)
    else:
        _render_text_only_body(slide, spec, body_top)

    _add_source_footer(slide, spec.source, spec.notes)


def _style_data_labels(plot, number_format: str, size=Pt(13), color=TEXT_DARK) -> None:
    """Turn on and style a plot's data labels, shared across chart types."""
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.number_format = number_format
    data_labels.number_format_is_linked = False
    data_labels.font.size = size
    data_labels.font.bold = False
    data_labels.font.color.rgb = color


def _style_category_axis(category_axis) -> None:
    """Shared tick-label styling for the bar and line charts' category axis."""
    category_axis.tick_labels.font.size = Pt(15)
    category_axis.tick_labels.font.name = FONT
    category_axis.tick_labels.font.color.rgb = TEXT_DARK


def _add_bar_chart(slide: Slide, spec: SlideSpec, top, height=None) -> None:
    """Horizontal bar chart for both ranking and comparison metrics.

    Ranking metrics are sorted ascending by rank (rank 1 on top);
    comparison metrics keep their given order. `height` defaults to
    filling the rest of the slide (see _chart_height).
    """
    metrics = spec.metrics
    if spec.chart_type == "ranking":
        # `m.rank is not None else 0` guards the LLM path: pydantic doesn't
        # require ranking metrics to carry a rank, so this avoids a
        # TypeError (None vs int) deep inside rendering.
        metrics = sorted(metrics, key=lambda m: m.rank if m.rank is not None else 0)

    # python-pptx bar charts plot the first category at the bottom, so
    # reverse to put rank 1 / highest value at the top.
    display_metrics = list(reversed(metrics))

    categories = [
        _truncate(f"{m.rank}. {m.label}" if spec.chart_type == "ranking" and m.rank else m.label, 40)
        for m in display_metrics
    ]
    values = [m.value for m in display_metrics]

    chart_data = CategoryChartData()
    chart_data.categories = categories
    series_name = "Rank (%)" if spec.chart_type == "ranking" else "Share (%)"
    chart_data.add_series(series_name, values)

    if height is None:
        height = _chart_height(top)
    width = Inches(11.5)
    left = Inches(0.9)

    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data)
    chart = graphic_frame.chart
    chart.has_legend = False
    chart.has_title = False

    plot = chart.plots[0]
    plot.gap_width = 60  # tighter bar spacing than PowerPoint's 150% default
    _style_data_labels(plot, number_format="0")

    # Every bar is the same color — no legend needed since category-axis
    # labels sit directly beside each bar.
    series = plot.series[0]
    for point in series.points:
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = BLUE_MEDIUM

    category_axis = chart.category_axis
    _style_category_axis(category_axis)
    category_axis.format.line.fill.background()

    value_axis = chart.value_axis
    value_axis.visible = False
    value_axis.has_major_gridlines = False


def _add_pie_chart(slide: Slide, spec: SlideSpec, top, height=None) -> None:
    """Pie chart for part-to-whole metrics (see planner._looks_like_part_to_whole).

    Each slice gets a distinct tint from PIE_PALETTE since there's no axis
    label to carry identity; a legend maps color to label instead.
    """
    # Clamp defensively: SlideSpec._pie_slice_count_in_range rejects an
    # out-of-range slice count on normal construction, but a slide built
    # via model_construct could still reach here.
    metrics = spec.metrics[:6]
    categories = [_truncate(m.label, 40) for m in metrics]
    values = [m.value for m in metrics]

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Share (%)", values)

    if height is None:
        height = _chart_height(top)
    width = Inches(8.0)
    left = Inches(0.9)

    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.PIE, left, top, width, height, chart_data)
    chart = graphic_frame.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(13)
    chart.legend.font.name = FONT
    chart.legend.font.color.rgb = TEXT_DARK

    plot = chart.plots[0]
    # Values arrive as plain numbers (57 meaning "57%"); Excel's built-in
    # "0%" format would multiply by 100 and misrender it as "5700%".
    _style_data_labels(plot, number_format='0"%"', size=Pt(12), color=RGBColor(0xFF, 0xFF, 0xFF))

    series = plot.series[0]
    for i, point in enumerate(series.points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = PIE_PALETTE[i % len(PIE_PALETTE)]


def _add_line_chart(slide: Slide, spec: SlideSpec, top, height=None) -> None:
    """Line chart for trend (time series) metrics, kept in their given
    (chronological) order. `height` defaults via _chart_height."""
    metrics = spec.metrics
    categories = [_truncate(m.label, 40) for m in metrics]
    values = [m.value for m in metrics]

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(spec.title, values)

    if height is None:
        height = _chart_height(top)
    width = Inches(11.5)
    left = Inches(0.9)

    graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, left, top, width, height, chart_data)
    chart = graphic_frame.chart
    chart.has_legend = False
    chart.has_title = False

    plot = chart.plots[0]
    _style_data_labels(plot, number_format="0.0")

    # One continuous primary-color line; every marker is a small dot except
    # the highlighted point, which is simply larger — same color throughout.
    series = plot.series[0]
    series.format.line.color.rgb = PRIMARY
    series.format.line.width = Pt(1.75)
    series.marker.format.fill.solid()
    series.marker.format.fill.fore_color.rgb = PRIMARY
    series.marker.format.line.color.rgb = PRIMARY
    series.marker.size = 6

    # zip() assumes points and metrics come back in the same order/count;
    # if a future python-pptx version breaks that, the highlight silently
    # disappears instead of crashing — logged so it's still discoverable.
    points = list(series.points)
    if len(points) != len(metrics):
        logger.warning(
            "trend chart point/metric count mismatch on %r: %d points vs %d metrics — "
            "highlight may land on the wrong point.",
            spec.title,
            len(points),
            len(metrics),
        )
    for point, metric in zip(points, metrics):
        if metric.highlight:
            point.marker.format.fill.solid()
            point.marker.format.fill.fore_color.rgb = PRIMARY
            point.marker.format.line.color.rgb = PRIMARY
            point.marker.size = 9

    category_axis = chart.category_axis
    _style_category_axis(category_axis)
    category_axis.format.line.color.rgb = GRIDLINE

    value_axis = chart.value_axis
    value_axis.tick_labels.font.size = Pt(12)
    value_axis.tick_labels.font.color.rgb = TEXT_MUTED
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = GRIDLINE
    value_axis.major_gridlines.format.line.width = Pt(0.25)
    value_axis.format.line.fill.background()


def _render_text_only_body(slide: Slide, spec: SlideSpec, top) -> None:
    """No chart_type -> text-only slide: key message lead-in + bullets.
    Bullets are capped to what fits the fixed-height box — anything beyond
    that is silently dropped so content never spills past the footer."""
    box_height = min(Inches(5.3), _chart_height(top))
    box = _add_textbox(slide, Inches(0.8), top, Inches(11.5), box_height)
    tf = box.text_frame
    first = True

    used_height = Emu(0)
    if spec.key_message:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = spec.key_message
        p.font.name = FONT
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PRIMARY
        p.space_after = Pt(18)
        first = False
        # Rough line-height estimate, just to budget remaining space for bullets.
        used_height = Pt(20 + 18)

    bullet_line_height = Pt(16 + 10)
    remaining_height = box_height - used_height
    max_bullets = max(1, int(remaining_height // bullet_line_height)) if spec.bullets else 0
    display_bullets = _cap_items(spec.bullets, max_bullets)

    for bullet in display_bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = f"•  {bullet}"
        p.font.name = FONT
        p.font.size = Pt(16)
        p.font.color.rgb = TEXT_DARK
        p.space_after = Pt(10)
        first = False

    if first:
        # Nothing to show — leave a visible placeholder so the slide isn't blank.
        p = tf.paragraphs[0]
        p.text = spec.title
        p.font.name = FONT
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_MUTED


# --- summary slide -------------------------------------------------------


def render_summary_slide(slide: Slide, spec: SlideSpec) -> None:
    """Closing recap slide: lead takeaway + numbered list of insight
    categories, each on its own row."""
    _add_slide_header(slide, spec)

    body_top = BODY_TOP + Inches(0.15)  # extra breathing room above the lead takeaway line

    if spec.key_message:
        msg_box = _add_textbox(slide, Inches(0.8), body_top, Inches(11.5), Inches(0.8))
        mp = msg_box.text_frame.paragraphs[0]
        mp.text = spec.key_message
        mp.font.name = FONT
        mp.font.size = Pt(19)
        mp.font.bold = True
        mp.font.color.rgb = PRIMARY
        items_top = body_top + Inches(0.95)
    else:
        items_top = body_top

    if spec.bullets:
        row_height = Inches(0.85)
        # Cap rows to what fits above the footer — anything beyond that is
        # silently dropped rather than running off the slide.
        available = FOOTER_TOP - Inches(0.15) - items_top
        max_rows = max(1, int(available // row_height))
        display_bullets = _cap_items(spec.bullets, max_rows)

        for i, bullet in enumerate(display_bullets):
            row_top = items_top + Emu(int(row_height * i))

            num_box = _add_textbox(slide, Inches(0.8), row_top, Inches(0.6), row_height)
            np = num_box.text_frame.paragraphs[0]
            np.text = f"{i + 1:02d}"
            np.font.name = FONT
            np.font.size = Pt(16)
            np.font.bold = True
            np.font.color.rgb = BLUE_MEDIUM

            text_box = _add_textbox(slide, Inches(1.5), row_top, Inches(10.6), row_height)
            tp = text_box.text_frame.paragraphs[0]
            tp.text = bullet
            tp.font.name = FONT
            tp.font.size = Pt(16)
            tp.font.color.rgb = TEXT_DARK

            if i < len(display_bullets) - 1:
                rule_top = row_top + row_height - Inches(0.12)
                rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), rule_top, Inches(11.2), Pt(0.75))
                rule.fill.solid()
                rule.fill.fore_color.rgb = GRIDLINE
                rule.line.fill.background()
                rule.shadow.inherit = False
    elif not spec.key_message:
        # Neither bullets nor key_message: fall back so the slide isn't blank.
        _render_text_only_body(slide, spec, body_top)

    _add_source_footer(slide, spec.source, spec.notes)


# --- KPI card slide ------------------------------------------------------


def render_kpi_slide(slide: Slide, spec: SlideSpec) -> None:
    """Headline-numbers slide: up to 4 KPI cards, each a flat white tile
    with a hairline border and the value in primary blue. Falls back to
    the text-only body if no kpis were provided."""
    _add_slide_header(slide, spec)

    body_top = BODY_TOP

    if not spec.kpis:
        _render_text_only_body(slide, spec, body_top)
        _add_source_footer(slide, spec.source, spec.notes)
        return

    if spec.key_message:
        _add_key_message(slide, spec.key_message, body_top)
        cards_top = body_top + Inches(0.8)
    else:
        cards_top = body_top + Inches(0.3)

    cards = spec.kpis[:4]

    n = len(cards)
    gutter = Inches(0.5)
    total_width = Inches(11.5)
    card_width = Emu(int((total_width - gutter * (n - 1)) / n))
    card_height = Inches(2.6)
    left0 = Inches(0.9)
    top_border_height = Pt(1.5)

    for i, kpi in enumerate(cards):
        left = Emu(int(left0 + i * (card_width + gutter)))

        _add_kpi_card(slide, left, cards_top, card_width, card_height, top_border_height)

        # KpiCard's max_length constraints are the primary guard; this is
        # a second line of defense for specs built via model_construct.
        value_text = _truncate(kpi.value, 20)
        label_text = _truncate(kpi.label, 40)

        value_box = _add_textbox(slide, left + Inches(0.15), cards_top + Inches(0.45), card_width - Inches(0.3), Inches(0.9))
        vp = value_box.text_frame.paragraphs[0]
        vp.text = value_text
        vp.font.name = FONT
        vp.font.size = Pt(30)
        vp.font.bold = False
        vp.font.color.rgb = PRIMARY
        vp.alignment = PP_ALIGN.CENTER

        label_box = _add_textbox(slide, left + Inches(0.15), cards_top + Inches(1.35), card_width - Inches(0.3), Inches(0.5))
        lp = label_box.text_frame.paragraphs[0]
        lp.text = label_text
        lp.font.name = FONT
        lp.font.size = Pt(13)
        lp.font.bold = True
        lp.font.color.rgb = TEXT_DARK
        lp.alignment = PP_ALIGN.CENTER

        if kpi.delta:
            delta_box = _add_textbox(slide, left + Inches(0.15), cards_top + Inches(1.85), card_width - Inches(0.3), Inches(0.5))
            dp = delta_box.text_frame.paragraphs[0]
            dp.text = _truncate(kpi.delta, 60)
            dp.font.name = FONT
            dp.font.size = Pt(11)
            dp.font.color.rgb = TEXT_MUTED
            dp.alignment = PP_ALIGN.CENTER

    _add_source_footer(slide, spec.source, spec.notes)
